"""
Inserts the TAGFC Data Flow slide into the existing PPTX after slide 10.
Run:  python insert_flow_slide.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import copy, os

PPTX_PATH = r"c:\Users\abhia\Desktop\counterfactual_basis_kernel-main\Report\TAGFC_Presentation.pptx"

# ── colours (same as main deck) ─────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x00, 0x33, 0x66)
MID_BLUE   = RGBColor(0x00, 0x5B, 0x96)
ORANGE     = RGBColor(0xE8, 0x6A, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF0, 0xF4, 0xF8)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
GREEN      = RGBColor(0x00, 0x77, 0x44)
LIGHT_BLUE = RGBColor(0xCC, 0xE5, 0xFF)
CREAM      = RGBColor(0xFF, 0xF5, 0xE6)

prs = Presentation(PPTX_PATH)
BLANK = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line_col=None, line_pt=0):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_col and line_pt:
        s.line.color.rgb = line_col; s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        bold=False, italic=False, size=13,
        color=DARK_GREY, align=PP_ALIGN.LEFT,
        font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.bold = bold; r.font.italic = italic
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.name = font
    return tb

def add_para(tf, text, bold=False, size=11, color=DARK_GREY,
             align=PP_ALIGN.LEFT, font="Calibri"):
    p = tf.add_paragraph(); p.alignment = align
    r = p.add_run(); r.text = text
    r.font.bold = bold; r.font.size = Pt(size)
    r.font.color.rgb = color; r.font.name = font

def banner(slide, title):
    rect(slide, 0, 0, 13.33, 1.1, fill=DARK_BLUE)
    txt(slide, title, 0.3, 0.1, 12.5, 0.85,
        bold=True, size=22, color=WHITE, align=PP_ALIGN.LEFT)
    rect(slide, 0, 1.1, 13.33, 0.05, fill=ORANGE)

def footer(slide, num, total=29):
    rect(slide, 0, 7.25, 13.33, 0.25, fill=DARK_BLUE)
    txt(slide, "TAGFC — Attention-Guided Frequency-Domain Counterfactual Explanations",
        0.3, 7.26, 11, 0.22, size=9, color=WHITE)
    txt(slide, f"{num} / {total}", 12.5, 7.26, 0.8, 0.22,
        size=9, color=WHITE, align=PP_ALIGN.RIGHT)

def step_box(slide, x, y, w, h, step_num, step_name, step_eq,
             bg_fill=LIGHT_BLUE, hdr_fill=MID_BLUE):
    rect(slide, x, y, w, h, fill=bg_fill,
         line_col=MID_BLUE, line_pt=1.2)
    rect(slide, x, y, w, 0.35, fill=hdr_fill)
    txt(slide, f"Step {step_num}", x+0.08, y+0.03, w-0.16, 0.3,
        bold=True, size=11, color=WHITE)
    txt(slide, step_name, x+0.08, y+0.38, w-0.16, 0.38,
        bold=True, size=12, color=DARK_BLUE)
    txt(slide, step_eq, x+0.08, y+0.76, w-0.16, h-0.85,
        size=9.5, color=DARK_GREY, font="Courier New", wrap=True)

def down_arrow(slide, x, cy, w=0.22, h=0.32, color=ORANGE):
    """Draw a simple downward arrow using a thin rect + triangle-ish shape."""
    # shaft
    shaft_w = 0.06
    rect(slide, x + (w-shaft_w)/2, cy, shaft_w, h*0.65, fill=color)
    # arrowhead (wider rect)
    rect(slide, x, cy + h*0.55, w, h*0.45, fill=color)

def right_arrow_label(slide, x, y, w, label, color=MID_BLUE):
    txt(slide, "→", x, y, 0.3, 0.35, bold=True, size=18,
        color=ORANGE, align=PP_ALIGN.CENTER)
    txt(slide, label, x+0.32, y+0.04, w-0.32, 0.3,
        size=10, color=color, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# BUILD THE NEW SLIDE
# ════════════════════════════════════════════════════════════════════════════
new_slide = prs.slides.add_slide(BLANK)
banner(new_slide, "TAGFC — How Each Step's Output Flows Into the Next Step")
footer(new_slide, "11*")   # placeholder number

# ── background ───────────────────────────────────────────────────────────────
rect(new_slide, 0.18, 1.18, 13.0, 5.92, fill=LIGHT_GREY)

# ── COLUMN LAYOUT ────────────────────────────────────────────────────────────
# Left column  (x≈0.25) : step boxes + arrows  (vertical flow)
# Right column (x≈5.2)  : "what output goes where" annotations
# Far right    (x≈9.8)  : feedback loop label

BOX_W = 4.7
BOX_H = 0.98
GAP   = 0.30       # gap between boxes (space for arrow)
X0    = 0.28
Y0    = 1.25

steps = [
    ("1", "Attention Rollout",
     "Â = 0.5·A + 0.5·I\nRollout = Â¹⊗…⊗Â^L\nOutput: s_t ∈ [0,1]^T"),
    ("2", "Haar DWT",
     "C_d = DWT(X[:,d]) ∈ R^L\nOutput: C ∈ R^{D×L}"),
    ("3", "Omega Weights",
     "G_wav = DWT(|∇X|)\nM = s_t × G_wav\nω = 1/normalize(M)"),
    ("4", "4-Term Objective",
     "L = L_flip + λ1·L_sparse(ω)\n    + λ2·L_cross + λ3·L_manifold(C)\nOutput: ∇L"),
    ("5", "Proximal GD + Soft-Thresh",
     "δ ← δ − η·∇L   [gradient step]\nδ_cd ← sign·max(|δ|−η·λ1·ω, 0)  [sparsify]\nOutput: δ  →  X* = iDWT(C + δ)"),
]

HDR_COLORS = [MID_BLUE, MID_BLUE, GREEN,
              RGBColor(0x66, 0x00, 0x99), ORANGE]
BG_COLORS  = [LIGHT_BLUE, LIGHT_BLUE, RGBColor(0xD4, 0xED, 0xDA),
              RGBColor(0xF0, 0xE0, 0xFF), CREAM]

for i, (num, name, eq) in enumerate(steps):
    y = Y0 + i*(BOX_H + GAP)
    step_box(new_slide, X0, y, BOX_W, BOX_H, num, name, eq,
             bg_fill=BG_COLORS[i], hdr_fill=HDR_COLORS[i])
    # down arrow between boxes
    if i < 4:
        # shaft
        ax = X0 + BOX_W/2 - 0.04
        ay = y + BOX_H
        rect(new_slide, ax, ay, 0.08, GAP*0.55, fill=ORANGE)
        # arrowhead
        rect(new_slide, ax-0.1, ay+GAP*0.48, 0.28, GAP*0.45, fill=ORANGE)

# ── Query X box (above Step 1) ───────────────────────────────────────────────
rect(new_slide, X0+0.3, 1.22, BOX_W-0.6, 0.38,
     fill=DARK_BLUE, line_col=DARK_BLUE, line_pt=0)
txt(new_slide, "Query X  (T×D multivariate time series)",
    X0+0.35, 1.23, BOX_W-0.7, 0.34,
    bold=True, size=12, color=WHITE, align=PP_ALIGN.CENTER)
# small arrow from query to step 1
rect(new_slide, X0+BOX_W/2-0.03, 1.60, 0.07, 0.28, fill=ORANGE)
rect(new_slide, X0+BOX_W/2-0.12, 1.82, 0.28, 0.22, fill=ORANGE)

# ── Counterfactual X* box (below Step 5) ─────────────────────────────────────
last_y = Y0 + 4*(BOX_H + GAP) + BOX_H
rect(new_slide, X0+0.3, last_y+0.12, BOX_W-0.6, 0.38,
     fill=GREEN, line_col=GREEN, line_pt=0)
txt(new_slide, "Counterfactual X*  =  iDWT(C + δ)",
    X0+0.35, last_y+0.13, BOX_W-0.7, 0.34,
    bold=True, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ── RIGHT PANEL: flow annotations ────────────────────────────────────────────
RX = 5.25   # right panel x start

# Panel background
rect(new_slide, RX-0.1, 1.18, 8.18, 5.92,
     fill=WHITE, line_col=MID_BLUE, line_pt=0.8)

# Header
rect(new_slide, RX-0.1, 1.18, 8.18, 0.38, fill=DARK_BLUE)
txt(new_slide, "What Each Step Produces  →  Where It Goes",
    RX, 1.21, 8.0, 0.33,
    bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)

connections = [
    # (y, from_color, from_label, arrow_label, to_label)
    (1.62,
     HDR_COLORS[0],
     "Step 1  →  s_t = [0.21, 0.24, 0.26, 0.29]",
     "Skips Step 2  |  Goes directly to Step 3",
     "Step 3 uses s_t to build joint importance M = s_t × G_wav"),

    (2.55,
     HDR_COLORS[1],
     "Step 2  →  C ∈ R^{2×4}  (wavelet coefficients)",
     "Used in Steps 4 and 5",
     "Step 4: X̃ = iDWT(C+δ) for loss computation\nStep 5: final X* = iDWT(C + δ_final)"),

    (3.48,
     HDR_COLORS[2],
     "Step 3  →  ω = [1.00, 2.04, 6.25, 2.44]",
     "Used in Steps 4 and 5",
     "Step 4: L_sparse = Σ ω·|δ|  (weighted cost)\nStep 5: threshold = η·λ₁·ω  (zeros out small δ)"),

    (4.41,
     HDR_COLORS[3],
     "Step 4  →  ∇L  (gradient of total loss)",
     "Used in Step 5 only",
     "Step 5: δ ← δ − η·∇L  (gradient step direction)"),

    (5.34,
     HDR_COLORS[4],
     "Step 5  →  δ  (sparse perturbation in wavelet space)",
     "Feeds BACK to Step 4  (loop 300 iterations)",
     "Step 4 recomputes loss with new X̃ = iDWT(C+δ)\nWhen P(Poor|X̃) > 0.5  →  EXIT loop  →  X* done"),
]

for (y, hdr_col, from_txt, arrow_txt, to_txt) in connections:
    # coloured "from" label
    rect(new_slide, RX, y, 7.9, 0.30, fill=hdr_col)
    txt(new_slide, from_txt, RX+0.1, y+0.03, 7.7, 0.26,
        bold=True, size=11, color=WHITE)
    # arrow label
    txt(new_slide, "  ↓  " + arrow_txt, RX+0.08, y+0.32, 7.8, 0.25,
        size=10.5, color=ORANGE, bold=True, italic=True)
    # "to" label
    txt(new_slide, to_txt, RX+0.2, y+0.56, 7.6, 0.40,
        size=10.5, color=DARK_GREY, wrap=True)

# ── feedback loop label (right margin) ───────────────────────────────────────
rect(new_slide, 12.85, 4.55, 0.42, 1.45,
     fill=DARK_BLUE, line_col=ORANGE, line_pt=1.5)
txt(new_slide, "↑\n↑\nloop\n↑\n↑",
    12.85, 4.55, 0.42, 1.45,
    size=9, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# One-sentence summary at bottom
rect(new_slide, 0.18, 7.05, 13.0, 0.22, fill=DARK_BLUE)
txt(new_slide,
    "Steps 1 & 2 run ONCE  ·  Step 3 runs ONCE  ·  Steps 4 ↔ 5 loop until convergence (max 300 iters)",
    0.22, 7.06, 12.9, 0.20,
    size=10.5, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

# ════════════════════════════════════════════════════════════════════════════
# INSERT slide at position 10 (0-indexed) = after current slide 10
# ════════════════════════════════════════════════════════════════════════════
# python-pptx appends slides; we need to move it to position 10
xml_slides = prs.slides._sldIdLst
# The new slide was appended at the end — move it to index 10 (after slide 10)
new_slide_elem = xml_slides[-1]
xml_slides.remove(new_slide_elem)
xml_slides.insert(10, new_slide_elem)   # 0-indexed: position 10 = after slide 10

print(f"Total slides after insertion: {len(prs.slides)}")
prs.save(PPTX_PATH)
print(f"Saved: {PPTX_PATH}")
print("New slide inserted at position 11 (after existing slide 10)")
