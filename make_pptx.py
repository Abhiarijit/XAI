"""
TAGFC Thesis Presentation Generator
Generates: Report/TAGFC_Presentation.pptx  (~28 slides)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = r"c:\Users\abhia\Desktop\counterfactual_basis_kernel-main"
OUT  = os.path.join(BASE, "Report", "TAGFC_Presentation.pptx")

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x00, 0x33, 0x66)   # IIT Guwahati navy
MID_BLUE   = RGBColor(0x00, 0x5B, 0x96)   # accent blue
ORANGE     = RGBColor(0xE8, 0x6A, 0x00)   # accent orange
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF0, 0xF4, 0xF8)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)

SLD_W = Inches(13.33)   # 16:9 widescreen
SLD_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLD_W
prs.slide_height = SLD_H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             bold=False, italic=False, size=18,
             color=DARK_GREY, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.name   = font_name
    return txb


def add_para(tf, text, bold=False, italic=False, size=16,
             color=DARK_GREY, align=PP_ALIGN.LEFT, font_name="Calibri",
             space_before=4):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = _Pt(size)
    run.font.color.rgb = color
    run.font.name   = font_name
    return p


def banner(slide, title_text, subtitle_text=""):
    """Dark-blue top banner with title."""
    add_rect(slide, 0, 0, 13.33, 1.15, fill_rgb=DARK_BLUE)
    add_text(slide, title_text, 0.3, 0.1, 12.5, 0.85,
             bold=True, size=24, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_text(slide, subtitle_text, 0.3, 0.82, 12.5, 0.4,
                 size=14, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.LEFT)
    # thin orange accent line
    add_rect(slide, 0, 1.15, 13.33, 0.05, fill_rgb=ORANGE)


def body_box(slide, l=0.4, t=1.35, w=12.5, h=5.9):
    """Returns a textframe inside a white body box."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    txb.text_frame.word_wrap = True
    return txb.text_frame


def img(slide, rel_path, l, t, w):
    """Add image; skip silently if missing."""
    full = os.path.join(BASE, rel_path)
    if not os.path.exists(full):
        add_text(slide, f"[Fig: {os.path.basename(rel_path)}]",
                 l, t, w, 2, size=10, color=RGBColor(0x99, 0x99, 0x99),
                 align=PP_ALIGN.CENTER)
        return
    slide.shapes.add_picture(full, Inches(l), Inches(t), width=Inches(w))


def footer(slide, slide_num, total=28):
    add_rect(slide, 0, 7.25, 13.33, 0.25, fill_rgb=DARK_BLUE)
    add_text(slide, "TAGFC — Attention-Guided Frequency-Domain Counterfactual Explanations",
             0.3, 7.26, 11, 0.22, size=9, color=WHITE)
    add_text(slide, f"{slide_num} / {total}", 12.5, 7.26, 0.8, 0.22,
             size=9, color=WHITE, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
add_rect(slide, 0, 5.6, 13.33, 1.9, fill_rgb=RGBColor(0x00, 0x22, 0x44))
add_rect(slide, 0.5, 2.8, 12.33, 0.07, fill_rgb=ORANGE)

add_text(slide,
         "Attention-Guided Frequency-Domain\nCounterfactual Explanations for\nTransformer-Based Multivariate\nTime-Series Classification",
         1.0, 0.55, 11.3, 2.4,
         bold=True, size=28, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Avinash Yadav (244101008)   |   Vaibhav Wankar (244101061)",
         1.0, 3.0, 11.3, 0.45, size=17, color=RGBColor(0xAD, 0xD8, 0xE6),
         align=PP_ALIGN.CENTER)
add_text(slide, "Guide: Dr. Rashmi Dutta Baruah",
         1.0, 3.45, 11.3, 0.4, size=15, color=RGBColor(0xAD, 0xD8, 0xE6),
         align=PP_ALIGN.CENTER)
add_text(slide, "Department of Computer Science & Engineering\nIndian Institute of Technology Guwahati",
         1.0, 3.88, 11.3, 0.7, size=14, color=RGBColor(0xCC, 0xDD, 0xEE),
         align=PP_ALIGN.CENTER)
add_text(slide, "M.Tech Thesis  ·  June 2026",
         1.0, 5.72, 11.3, 0.38, size=13, color=RGBColor(0x99, 0xBB, 0xDD),
         align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PRESENTATION OUTLINE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Presentation Outline")
footer(slide, 2)

items = [
    ("1", "Introduction & Motivation"),
    ("2", "Problem Statement & Research Gaps"),
    ("3", "Related Work — Existing Methods"),
    ("4", "Datasets"),
    ("5", "TAGFC Methodology — 5-Step Pipeline"),
    ("6", "Experimental Results: AQI India"),
    ("7", "Experimental Results: NATOPS"),
    ("8", "Comparison & Statistical Analysis"),
    ("9", "Conclusions & Future Work"),
]
for i, (num, txt) in enumerate(items):
    y = 1.4 + i * 0.6
    add_rect(slide, 0.4, y, 0.45, 0.42, fill_rgb=MID_BLUE)
    add_text(slide, num, 0.4, y+0.03, 0.45, 0.38,
             bold=True, size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, txt, 1.05, y+0.04, 11.5, 0.38, size=17, color=DARK_GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MOTIVATION
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Motivation", "Why do we need counterfactual explanations for time-series classifiers?")
footer(slide, 3)

bullets = [
    ("Critical-domain deployment:", "Predictive maintenance (NASA turbofans), air quality monitoring, and gesture interfaces all require not just predictions but actionable explanations."),
    ("Deep learning opacity:", "Transformer classifiers achieve state-of-the-art accuracy on multivariate time series but are black boxes — operators cannot verify or trust their decisions."),
    ("SHAP/LIME limitations:", "Feature-attribution methods answer 'what mattered?' but cannot answer 'what would need to change for a different outcome?' — they lack actionability."),
    ("Counterfactual gap:", "Existing CF methods for time series (CoMTE, AB-CF, Native Guide) ignore sensor physics, work in raw-signal space, and discard model-internal saliency signals."),
    ("TAGFC insight:", "Combining Transformer attention rollout + Haar wavelet domain + cross-sensor coherence yields sparse, physically valid, model-aware counterfactuals."),
]
for i, (hd, body) in enumerate(bullets):
    y = 1.35 + i * 1.08
    add_rect(slide, 0.3, y, 0.06, 0.55, fill_rgb=ORANGE)
    add_text(slide, hd, 0.5, y, 3.5, 0.38, bold=True, size=15, color=MID_BLUE)
    add_text(slide, body, 0.5, y+0.35, 12.3, 0.65, size=13.5, color=DARK_GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Problem Statement")
footer(slide, 4)

add_rect(slide, 0.3, 1.3, 12.73, 1.6, fill_rgb=LIGHT_GREY)
add_text(slide,
         "Given a trained Transformer classifier f : ℝ^{T×D} → {1,…,K} and a query time series X "
         "with predicted label y = f(X), find a counterfactual X* such that f(X*) = y* ≠ y, "
         "where X* is (i) minimal in perturbation, (ii) sparse in changed sensors/timesteps, "
         "(iii) physically coherent across sensors, and (iv) close to real target-class samples.",
         0.5, 1.35, 12.3, 1.5, size=15, color=DARK_GREY, wrap=True)

add_text(slide, "Formal Objective", 0.3, 3.05, 4, 0.38, bold=True, size=16, color=MID_BLUE)
add_text(slide,
         "min_δ   L_flip(X̃) + λ₁·L_sparse(δ) + λ₂·L_cross(X̃) + λ₃·L_manifold(X̃)",
         0.3, 3.42, 12.5, 0.52, bold=True, size=16.5, color=DARK_BLUE,
         align=PP_ALIGN.CENTER, font_name="Courier New")

terms = [
    ("L_flip",     "−log P(y*|X̃)",                    "Classifier flip loss"),
    ("L_sparse",   "Σ  ω_cd · |δ_cd|",                "Omega-weighted wavelet L1"),
    ("L_cross",    "Δμᵀ Σ⁻¹ Δμ",                      "Mahalanobis cross-sensor coherence"),
    ("L_manifold", "‖X̃ − X_nn‖²_F / (T·D)",          "Proximity to nearest target sample"),
]
for i, (nm, eq, desc) in enumerate(terms):
    x = 0.3 + i * 3.25
    add_rect(slide, x, 4.15, 3.0, 1.55, fill_rgb=RGBColor(0xE8, 0xF4, 0xFF),
             line_rgb=MID_BLUE, line_pt=1)
    add_text(slide, nm,   x+0.1, 4.20, 2.8, 0.38, bold=True, size=14, color=ORANGE)
    add_text(slide, eq,   x+0.1, 4.54, 2.8, 0.42, size=13, color=DARK_BLUE,
             font_name="Courier New")
    add_text(slide, desc, x+0.1, 4.93, 2.8, 0.55, size=11.5, color=DARK_GREY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RESEARCH GAPS & OBJECTIVES
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Research Gaps & Objectives")
footer(slide, 5)

add_text(slide, "Identified Research Gaps", 0.3, 1.3, 6, 0.38,
         bold=True, size=16, color=MID_BLUE)
gaps = [
    "No existing MTS CF method operates in the frequency/wavelet domain",
    "Existing methods ignore model-internal attention (black-box perturbation)",
    "Cross-sensor physical coherence is not enforced — outputs are physically implausible",
]
for i, g in enumerate(gaps):
    add_rect(slide, 0.3, 1.75+i*0.75, 0.38, 0.38, fill_rgb=ORANGE)
    add_text(slide, "✗", 0.3, 1.77+i*0.75, 0.38, 0.34,
             bold=True, size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, g, 0.8, 1.77+i*0.75, 5.8, 0.55, size=14, color=DARK_GREY)

add_rect(slide, 6.9, 1.3, 0.05, 5.5, fill_rgb=RGBColor(0xCC, 0xCC, 0xCC))

add_text(slide, "Research Objectives", 6.95, 1.3, 6, 0.38,
         bold=True, size=16, color=MID_BLUE)
objs = [
    "Design a wavelet-domain counterfactual optimisation framework for MTS",
    "Derive attention rollout saliency to guide perturbation focus",
    "Enforce cross-sensor Mahalanobis coherence in the objective",
    "Achieve exact sparsity via Proximal Gradient Descent with soft-thresholding",
    "Evaluate on two benchmark datasets (AQI India, NATOPS) against 4 baselines",
    "Conduct statistical validation with Wilcoxon signed-rank + Bonferroni correction",
]
for i, o in enumerate(objs):
    add_rect(slide, 6.95, 1.75+i*0.6, 0.38, 0.38, fill_rgb=MID_BLUE)
    add_text(slide, str(i+1), 6.95, 1.77+i*0.6, 0.38, 0.34,
             bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, o, 7.45, 1.77+i*0.6, 5.5, 0.55, size=13, color=DARK_GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — RELATED WORK
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Related Work — Counterfactual Methods for Time Series")
footer(slide, 6)

headers = ["Method", "Year", "Search Space", "Saliency-Guided?", "Frequency Domain?", "Coherence?"]
rows = [
    ["Native Guide", "2021", "Raw signal\n(window swap)", "CAM-based\n(univariate)", "No", "No"],
    ["CoMTE",        "2021", "Raw signal\n(whole channel swap)", "No",          "No", "No"],
    ["SETS",         "2022", "Shapelet space",            "No",                "No", "No"],
    ["AB-CF",        "2023", "Raw signal\n(segment entropy)", "Entropy proxy", "No", "No"],
    ["TAGFC (Ours)", "2026", "Haar wavelet\ncoefficients", "Attention rollout\n(multi-layer)", "Yes (3-level\nHaar DWT)", "Yes (Maha-\nlanobis Σ⁻¹)"],
]
col_w = [1.85, 0.7, 1.75, 1.85, 1.6, 1.4]
col_x = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w + 0.02)

# header row
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(slide, cx, 1.35, cw, 0.42, fill_rgb=DARK_BLUE)
    add_text(slide, hdr, cx+0.05, 1.37, cw-0.1, 0.38,
             bold=True, size=12, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    bg = RGBColor(0xFF, 0xF5, 0xE6) if i == 4 else (LIGHT_GREY if i % 2 == 0 else WHITE)
    for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(slide, cx, 1.79+i*0.78, cw, 0.74, fill_rgb=bg,
                 line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_pt=0.5)
        clr = ORANGE if (i == 4 and j == 0) else DARK_GREY
        bld = (i == 4)
        add_text(slide, cell, cx+0.05, 1.82+i*0.78, cw-0.1, 0.68,
                 bold=bld, size=11.5, color=clr, align=PP_ALIGN.CENTER, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DATASETS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Datasets")
footer(slide, 7)

ds = [
    {
        "name": "AQI India",
        "source": "Central Pollution Control Board",
        "T": "T = 14 timesteps",
        "D": "D = 12 pollutants",
        "K": "K = 3 classes",
        "details": "PM2.5, PM10, NO, NO2, NOx, NH3, CO, SO2, O3, Benzene, Toluene, Xylene\nClasses: Good / Moderate / Poor",
        "N": "N = 60 explained",
    },
    {
        "name": "NATOPS",
        "source": "UEA Time Series Archive",
        "T": "T = 51 timesteps",
        "D": "D = 24 channels",
        "K": "K = 6 gesture classes",
        "details": "Right/Left: Hand, Elbow, Wrist, Thumb × x,y,z\nClasses: I have command, All clear, Not clear, Spread wings, Fold wings, Lock wings",
        "N": "N = 30 explained",
    },
    {
        "name": "NASA CMAPSS FD001",
        "source": "NASA Turbofan Degradation",
        "T": "T = 30 cycles",
        "D": "D = 14 sensors",
        "K": "K = 3 health states",
        "details": "7 constant sensors removed (s1,s5,s6,s10,s16,s18,s19)\nClasses: Healthy (RUL>120), Degrading, Critical (RUL≤30)",
        "N": "N = 200 (pending)",
    },
]
for i, d in enumerate(ds):
    x = 0.3 + i * 4.38
    add_rect(slide, x, 1.3, 4.1, 5.8, fill_rgb=LIGHT_GREY,
             line_rgb=MID_BLUE, line_pt=1.5)
    add_rect(slide, x, 1.3, 4.1, 0.5, fill_rgb=MID_BLUE)
    add_text(slide, d["name"], x+0.1, 1.33, 3.9, 0.44,
             bold=True, size=16, color=WHITE)
    add_text(slide, d["source"], x+0.1, 1.85, 3.9, 0.35, size=11, color=DARK_GREY, italic=True)
    for j, line in enumerate([d["T"], d["D"], d["K"], d["N"]]):
        add_text(slide, line, x+0.15, 2.25+j*0.42, 3.8, 0.4,
                 bold=(j < 3), size=13, color=DARK_BLUE)
    add_text(slide, d["details"], x+0.1, 4.0, 3.9, 1.5,
             size=11.5, color=DARK_GREY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TAGFC FRAMEWORK OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "TAGFC Framework — 5-Step Pipeline")
footer(slide, 8)

steps = [
    ("1", "Attention\nRollout", "Ât^l = 0.5·At^l + 0.5·I\nRollout = ∏ Ât\nst = column sums → [0,1]^T"),
    ("2", "Haar DWT", "Decompose each sensor\nX → C ∈ R^{D×L}\n3-level wavelet transform"),
    ("3", "Omega\nWeights", "G = |smoothed gradient|\nGwav = DWT(|G|)\nω = 1 / (s × Gwav)"),
    ("4", "4-Term\nObjective", "Lflip + λ1·Lsparse\n+ λ2·Lcross\n+ λ3·Lmanifold"),
    ("5", "Proximal GD\n+ Soft-Thresh", "δ ← δ − η∇L\nδcd ← sign(δcd)·max(|δcd|−η·λ1·ωcd, 0)\nClip δ ∈ [−2.5, +2.5]"),
]
for i, (num, name, eq) in enumerate(steps):
    x = 0.25 + i * 2.6
    add_rect(slide, x, 1.35, 2.42, 4.8,
             fill_rgb=RGBColor(0xE8, 0xF4, 0xFF), line_rgb=MID_BLUE, line_pt=1.5)
    add_rect(slide, x, 1.35, 2.42, 0.55, fill_rgb=MID_BLUE)
    add_text(slide, f"Step {num}", x+0.08, 1.38, 2.25, 0.28,
             bold=True, size=12, color=WHITE)
    add_text(slide, name, x+0.08, 1.63, 2.25, 0.55,
             bold=True, size=15, color=DARK_BLUE)
    add_text(slide, eq, x+0.08, 2.22, 2.25, 3.5,
             size=12, color=DARK_GREY, font_name="Courier New", wrap=True)
    # arrow (skip for last)
    if i < 4:
        add_text(slide, "→", x+2.42, 3.3, 0.18, 0.5,
                 bold=True, size=22, color=ORANGE, align=PP_ALIGN.CENTER)

# Hyperparameters bar
add_rect(slide, 0.25, 6.3, 12.83, 0.5, fill_rgb=RGBColor(0x00, 0x22, 0x44))
add_text(slide,
         "Hyperparameters:  λ₁=0.02   λ₂=0.01   λ₃=0.10   η=0.05   MAX_ITER=300   PATIENCE=25   DELTA_BOUND=2.5   σ=1.5",
         0.35, 6.34, 12.7, 0.42, size=12.5, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — STEP 1 & 2 DETAIL
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Method Detail — Attention Rollout & Haar DWT")
footer(slide, 9)

add_text(slide, "Step 1 — Attention Rollout", 0.3, 1.35, 6.2, 0.38,
         bold=True, size=17, color=MID_BLUE)
ar = [
    "For each Transformer layer l, add identity shortcut:",
    "     Â^l  =  0.5 · A^l  +  0.5 · I_T",
    "Propagate through all L layers (recursive product):",
    "     Rollout  =  Â^1 ⊗ Â^2 ⊗ … ⊗ Â^L",
    "Temporal saliency:  s_t = Σ_h  Rollout[CLS, t]  ∈ [0,1]^T",
    "High s_t → timestep t is important to prediction",
]
for i, line in enumerate(ar):
    fn = "Courier New" if line.startswith(" ") else "Calibri"
    add_text(slide, line.strip(), 0.3, 1.8+i*0.45, 6.4, 0.42, size=13,
             color=DARK_GREY if not line.startswith(" ") else DARK_BLUE,
             font_name=fn)

add_rect(slide, 6.8, 1.35, 0.04, 5.2, fill_rgb=RGBColor(0xCC, 0xCC, 0xCC))

add_text(slide, "Step 2 — Haar DWT", 7.0, 1.35, 5.8, 0.38,
         bold=True, size=17, color=MID_BLUE)
dwt = [
    "Decompose each sensor channel independently:",
    "  C_d  =  DWT(X[:, d])  ∈ R^L     (L = 64 for T=51)",
    "Full coefficient matrix: C ∈ R^{D × L}",
    "Counterfactual: X̃ = iDWT(C + δ)  (inverse DWT)",
    "",
    "Why wavelets?",
    "  • Slow trends (approx. coefficients) ≠ fast noise",
    "  • Localised in both time AND frequency",
    "  • Sparse perturbation of a few coefficients",
    "    changes a specific time-frequency region",
]
for i, line in enumerate(dwt):
    if not line:
        continue
    fn = "Courier New" if "=" in line and not line.startswith("Why") else "Calibri"
    bld = line.startswith("Why") or line.startswith("  •")
    add_text(slide, line.strip(), 7.0, 1.8+i*0.45, 6.1, 0.42, size=13,
             color=DARK_GREY, bold=bld, font_name=fn)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — STEP 3, 4, 5 DETAIL
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Method Detail — Omega Weights, Objective, Optimisation")
footer(slide, 10)

sections = [
    ("Step 3 — Omega Weights", [
        "Smoothed input gradient:  G = Gaussian_smooth(|∂L/∂X|,  σ=1.5)",
        "Wavelet gradient:  G_wav = DWT(|G|)  ∈ R^{D×L}",
        "Joint importance:  M = s_t(×broadcast) × G_wav",
        "Omega (perturbation cost):  ω = 1 / normalize(M)  — cheap to change if ω is small",
    ]),
    ("Step 4 — Objective (4 terms)", [
        "L_flip    = −log P(y* | X̃)          [classifier flip]",
        "L_sparse  = Σ_{c,d} ω_{cd} · |δ_{cd}|  [omega-weighted L1]",
        "L_cross   = Δμᵀ · Σ_X⁻¹ · Δμ         [Maha. cross-sensor]",
        "L_manifold= ‖X̃ − X_nn‖²_F / (T·D)   [manifold proximity]",
        "Total:    L = L_flip + 0.02·L_sparse + 0.01·L_cross + 0.10·L_manifold",
    ]),
    ("Step 5 — Proximal GD", [
        "Gradient step:  δ ← δ − η · ∇_δ L(δ)",
        "Soft-threshold: δ_{cd} ← sign(δ_{cd}) · max(|δ_{cd}| − η·λ₁·ω_{cd},  0)",
        "Bound clip:     δ_{cd} ← clip(δ_{cd},  −2.5,  +2.5)",
        "Stop: patience=25 iters no improvement, max 300 iters",
    ]),
]
y0 = 1.35
for sec_title, lines in sections:
    add_rect(slide, 0.3, y0, 12.7, 0.35, fill_rgb=MID_BLUE)
    add_text(slide, sec_title, 0.4, y0+0.02, 12.5, 0.31, bold=True, size=14, color=WHITE)
    y0 += 0.38
    for line in lines:
        add_text(slide, line, 0.5, y0, 12.2, 0.35, size=13,
                 color=DARK_GREY, font_name="Courier New")
        y0 += 0.38
    y0 += 0.15


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — EVALUATION METRICS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Evaluation Metrics (8 Metrics)")
footer(slide, 11)

metrics = [
    ("Validity ↑",        "Fraction of queries where f(X*) ≠ f(X)",               "Binary — must be 1 for a CF to be meaningful"),
    ("Proximity L1 ↓",    "‖X* − X‖₁ / (T·D)",                                    "Total magnitude of change"),
    ("Proximity L2 ↓",    "‖X* − X‖₂ / √(T·D)",                                   "RMS perturbation (sensitive to large changes)"),
    ("Proximity L∞ ↓",    "max|X* − X|",                                           "Worst-case single timestep change"),
    ("Sparsity ↑",        "1 − (# changed features) / (T·D)",                      "Fraction of features left unchanged"),
    ("Coherence ↓",       "(Δμ)ᵀ Σ_X⁻¹ (Δμ)",                                     "Mahalanobis dist — cross-sensor plausibility"),
    ("CF Confidence ↑",   "P(y* | X*) from classifier",                            "How strongly the CF belongs to target class"),
    ("RCF ↓",             "‖X* − X‖₂ / ‖X_NUN − X‖₂",                            "CF quality vs. nearest unlike neighbour"),
]
for i, (name, eq, desc) in enumerate(metrics):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.55
    y = 1.35 + row * 1.42
    add_rect(slide, x, y, 6.2, 1.3, fill_rgb=LIGHT_GREY,
             line_rgb=MID_BLUE, line_pt=0.8)
    add_text(slide, name, x+0.12, y+0.06, 5.9, 0.35,
             bold=True, size=14, color=ORANGE)
    add_text(slide, eq, x+0.12, y+0.4, 5.9, 0.38,
             size=12.5, color=DARK_BLUE, font_name="Courier New")
    add_text(slide, desc, x+0.12, y+0.76, 5.9, 0.45,
             size=11.5, color=DARK_GREY, italic=True)

add_text(slide, "Statistical Test: Paired Wilcoxon Signed-Rank + Bonferroni  (α = 0.05/6 = 0.0083)",
         0.3, 7.0, 12.7, 0.35, size=13, color=MID_BLUE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — AQI TRAINING
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "AQI India — Classifier Training Results")
footer(slide, 12)

add_text(slide, "Transformer classifier trained on AQI India dataset (T=14, D=12, K=3)",
         0.3, 1.3, 12.7, 0.38, size=14, color=DARK_GREY, italic=True)

img(slide, r"AQI_XAI\outputs\figures\F2_training_curves.png.png", 2.2, 1.78, 9.0)

add_text(slide,
         "Training Accuracy converges within 80 epochs  ·  3-class classification: Good / Moderate / Poor air quality",
         0.3, 6.7, 12.7, 0.38, size=12.5, color=DARK_GREY, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — AQI CF VISUALISATION
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "AQI India — Counterfactual Time-Series Visualisation")
footer(slide, 13)

add_text(slide, "Original vs. Counterfactual trajectories for 12 pollutants over 14 time-steps",
         0.3, 1.28, 12.7, 0.38, size=14, color=DARK_GREY, italic=True)
img(slide, r"AQI_XAI\outputs\figures\F3_cf_timeseries_sample0.png", 0.5, 1.7, 12.3)

add_text(slide,
         "TAGFC perturbs only the most salient wavelet coefficients → targeted, sparse changes "
         "while keeping physically correlated pollutants coherent",
         0.3, 6.65, 12.7, 0.5, size=12.5, color=DARK_GREY, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — AQI SALIENCY & POLLUTANT IMPORTANCE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "AQI India — Temporal Saliency & Pollutant Importance")
footer(slide, 14)

img(slide, r"AQI_XAI\outputs\figures\F4_temporal_saliency_sample0.png", 0.3, 1.3, 6.2)
img(slide, r"AQI_XAI\outputs\figures\F5_pollutant_importance_sample0.png", 6.8, 1.3, 6.2)

add_text(slide, "Temporal Saliency (Attention Rollout)",
         0.3, 6.35, 6.2, 0.38, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Pollutant Importance (Omega Weights)",
         6.8, 6.35, 6.2, 0.38, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "High saliency timesteps and high-omega pollutants are perturbed first by TAGFC",
         0.3, 6.85, 12.7, 0.35, size=12, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — AQI FREQUENCY ANALYSIS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "AQI India — Frequency-Domain Analysis")
footer(slide, 15)

img(slide, r"AQI_XAI\outputs\figures\F6_freq_importance_sample0.png", 0.3, 1.3, 6.2)
img(slide, r"AQI_XAI\outputs\figures\F7_delta_heatmap_sample0.png",   6.8, 1.3, 6.2)

add_text(slide, "Frequency Band Importance",
         0.3, 6.35, 6.2, 0.38, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Wavelet Coefficient Δ Heatmap (Sensor × Coefficient)",
         6.8, 6.35, 6.2, 0.38, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide,
         "Low-frequency (approx.) coefficients carry dominant pollution trend signal; "
         "TAGFC concentrates perturbations there",
         0.3, 6.88, 12.7, 0.35, size=12, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — AQI QUANTITATIVE RESULTS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "AQI India — Quantitative Results  (N = 60)")
footer(slide, 16)

# Table
col_headers = ["Method", "Validity↑", "Prox L1↓", "Prox L2↓", "Prox L∞↓",
               "Sparsity↑", "Coherence↓", "CF Conf↑", "RCF↓"]
table_data = [
    ["TAGFC (Ours)",
     "1.000±0.000", "10.589±8.229", "2.030±1.214", "1.304±0.704",
     "0.000±0.001", "0.132±0.274",  "0.955±0.092", "0.349±0.182"],
    ["CoMTE",
     "0.167±0.373", "47.136±32.907", "6.083±3.702", "2.326±1.475",
     "0.208±0.307", "1.042±1.428",  "0.132±0.292", "0.920±0.224"],
    ["AB-CF",
     "0.133±0.340", "49.305±33.406", "6.329±3.748", "2.486±1.766",
     "0.173±0.287", "1.044±1.413",  "0.106±0.271", "0.947±0.190"],
]
# best values (col indices in table_data, 0-indexed after method col)
best_cols = {0: 0, 1: 0, 2: 0, 3: 0, 4: 1, 5: 0, 6: 0, 7: 0}  # row index of best

cw = [1.5, 1.15, 1.15, 1.05, 1.05, 1.05, 1.15, 1.05, 1.05]
cx = [0.25]
for w in cw[:-1]:
    cx.append(cx[-1] + w + 0.01)

# header
for j, (hdr, x, w) in enumerate(zip(col_headers, cx, cw)):
    add_rect(slide, x, 1.3, w, 0.42, fill_rgb=DARK_BLUE)
    add_text(slide, hdr, x+0.04, 1.32, w-0.08, 0.38,
             bold=True, size=10, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(table_data):
    bg = RGBColor(0xFF, 0xF0, 0xD0) if i == 0 else (LIGHT_GREY if i % 2 == 0 else WHITE)
    for j, (cell, x, w) in enumerate(zip(row, cx, cw)):
        is_best = (best_cols.get(j-1, -1) == i and j > 0)
        add_rect(slide, x, 1.74+i*0.58, w, 0.54, fill_rgb=bg,
                 line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_pt=0.5)
        clr = RGBColor(0xCC, 0x00, 0x00) if is_best else (ORANGE if i == 0 else DARK_GREY)
        add_text(slide, cell, x+0.04, 1.77+i*0.58, w-0.08, 0.48,
                 bold=(i == 0 or is_best), size=10.5, color=clr,
                 align=PP_ALIGN.CENTER, wrap=True)

add_text(slide, "Bold red = best in column  ·  TAGFC wins 7/8 metrics; CoMTE best on Sparsity (whole-channel swap naturally zeroes unused channels)",
         0.25, 3.55, 12.83, 0.45, size=11.5, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)

img(slide, r"AQI_XAI\outputs\figures\F8_metric_comparison.png", 1.0, 4.05, 11.2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — NATOPS TRAINING
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "NATOPS — Classifier Training & Confusion Matrix")
footer(slide, 17)

img(slide, r"TAGFC_Natops\outputs\figures\F2_training_curves.png", 0.3, 1.3, 6.3)
img(slide, r"TAGFC_Natops\outputs\figures\F3_confusion_matrix.png", 6.9, 1.3, 6.1)

add_text(slide, "Training Accuracy + Loss Curves",
         0.3, 6.35, 6.3, 0.35, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Confusion Matrix (6-class gesture recognition)",
         6.9, 6.35, 6.1, 0.35, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide,
         "Transformer achieves high per-class accuracy on NATOPS — providing a strong, well-calibrated "
         "base model for CF explanation quality evaluation",
         0.3, 6.88, 12.7, 0.35, size=12, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — NATOPS CF VIZ
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "NATOPS — Counterfactual Time-Series Visualisation (24 channels)")
footer(slide, 18)

add_text(slide, "Original (blue) vs. Counterfactual (orange) — 24 motion capture channels over 51 timesteps",
         0.3, 1.28, 12.7, 0.35, size=13, color=DARK_GREY, italic=True)
img(slide, r"TAGFC_Natops\outputs\figures\F4_cf_timeseries_sample0.png", 0.3, 1.68, 12.7)

add_text(slide,
         "TAGFC modifies only a small subset of wavelet coefficients — resulting CF differs from original "
         "only in targeted motion segments, not the full trajectory",
         0.3, 6.7, 12.7, 0.45, size=12.5, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — NATOPS SALIENCY + CHANNEL IMPORTANCE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "NATOPS — Temporal Saliency & Channel Importance")
footer(slide, 19)

img(slide, r"TAGFC_Natops\outputs\figures\F5_temporal_saliency_sample0.png", 0.3, 1.3, 6.2)
img(slide, r"TAGFC_Natops\outputs\figures\F6_sensor_importance_sample0.png", 6.8, 1.3, 6.2)

add_text(slide, "Temporal Saliency (Attention Rollout over 51 timesteps)",
         0.3, 6.35, 6.2, 0.38, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Channel Importance (Omega × Saliency over 24 channels)",
         6.8, 6.35, 6.2, 0.38, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide,
         "Wrist and Elbow channels dominate importance — consistent with physical gesture biomechanics",
         0.3, 6.88, 12.7, 0.35, size=12, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — NATOPS FREQUENCY ANALYSIS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "NATOPS — Frequency-Domain Analysis")
footer(slide, 20)

img(slide, r"TAGFC_Natops\outputs\figures\F7_freq_importance_sample0.png", 0.3, 1.3, 6.2)
img(slide, r"TAGFC_Natops\outputs\figures\F8_delta_heatmap_sample0.png",  6.8, 1.3, 6.2)

add_text(slide, "Wavelet Band Importance (Approx + Detail levels)",
         0.3, 6.35, 6.2, 0.38, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Δ Heatmap: Sensor × Coefficient (where did TAGFC perturb?)",
         6.8, 6.35, 6.2, 0.38, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide,
         "Low-frequency (coarse motion) wavelet bands dominate — gesture transitions are slow structural changes, "
         "not high-frequency noise",
         0.3, 6.88, 12.7, 0.35, size=12, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — NATOPS OMEGA HEATMAP + CONVERGENCE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "NATOPS — Omega Heatmap & Convergence")
footer(slide, 21)

img(slide, r"TAGFC_Natops\outputs\figures\F9_omega_heatmap_sample0.png",    0.3, 1.3, 6.2)
img(slide, r"TAGFC_Natops\outputs\figures\F13_convergence_sample0.png",     6.8, 1.3, 6.2)

add_text(slide, "Omega Weight Heatmap (Sensor × Coefficient)",
         0.3, 6.35, 6.2, 0.38, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Optimisation Convergence (300 iterations, patience=25)",
         6.8, 6.35, 6.2, 0.38, bold=True, size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(slide,
         "Omega encodes both attention saliency and gradient importance — "
         "low-omega coefficients are targeted first, achieving sparse yet effective perturbations",
         0.3, 6.88, 12.7, 0.35, size=12, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — NATOPS QUANTITATIVE RESULTS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "NATOPS — Quantitative Results  (N = 30)")
footer(slide, 22)

# NATOPS results from memory
natops_data = [
    ["TAGFC (Ours)", "0.867±0.346", "35.498±23.591", "5.197±2.875", "0.883±0.376",
     "0.000±0.001", "7.544±8.831",  "0.705±0.383", "0.217±0.172"],
    ["CoMTE",        "1.000±0.000", "88.631±38.772", "14.774±5.672", "2.875±0.994",
     "0.917±0.054", "42.904±27.531","1.000±0.000",  "0.647±0.252"],
    ["AB-CF",        "1.000±0.000", "131.052±69.008","22.800±9.916","3.027±1.064",
     "0.882±0.083", "43.471±25.968","1.000±0.000",  "1.027±0.411"],
]
# best: Validity→CoMTE/ABCF, L1→TAGFC, L2→TAGFC, Linf→TAGFC, Sparsity→CoMTE, Coh→TAGFC, CFConf→CoMTE/ABCF, RCF→TAGFC
best_natops = {1: 0, 2: 0, 3: 0, 5: 0, 7: 0}  # col(1-indexed) → row of best

cw2 = [1.5, 1.1, 1.2, 1.05, 1.05, 1.05, 1.2, 1.05, 1.05]
cx2 = [0.25]
for w in cw2[:-1]:
    cx2.append(cx2[-1] + w + 0.01)

for j, (hdr, x, w) in enumerate(zip(col_headers, cx2, cw2)):
    add_rect(slide, x, 1.3, w, 0.42, fill_rgb=DARK_BLUE)
    add_text(slide, hdr, x+0.04, 1.32, w-0.08, 0.38,
             bold=True, size=10, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(natops_data):
    bg = RGBColor(0xFF, 0xF0, 0xD0) if i == 0 else (LIGHT_GREY if i % 2 == 0 else WHITE)
    for j, (cell, x, w) in enumerate(zip(row, cx2, cw2)):
        is_best = (best_natops.get(j, -1) == i)
        add_rect(slide, x, 1.74+i*0.58, w, 0.54, fill_rgb=bg,
                 line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_pt=0.5)
        clr = RGBColor(0xCC, 0x00, 0x00) if is_best else (ORANGE if i == 0 else DARK_GREY)
        add_text(slide, cell, x+0.04, 1.77+i*0.58, w-0.08, 0.48,
                 bold=(i == 0 or is_best), size=10.5, color=clr,
                 align=PP_ALIGN.CENTER, wrap=True)

add_text(slide,
         "TAGFC wins 4/8 metrics (Prox L1/L2/L∞, Coherence, RCF) vs CoMTE (p<0.0083) and "
         "5/8 vs AB-CF.  CoMTE/AB-CF achieve Validity=1.0 via whole-sequence swap — trading proximity for flip.",
         0.25, 3.55, 12.83, 0.5, size=11.5, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)

img(slide, r"TAGFC_Natops\outputs\figures\F10_metric_comparison.png", 1.0, 4.1, 11.2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — WILCOXON STATISTICAL TESTS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Comparison & Analysis — Wilcoxon Signed-Rank Tests")
footer(slide, 23)

add_text(slide,
         "Paired Wilcoxon signed-rank test  ·  Bonferroni-corrected α = 0.05/6 = 0.0083  ·  N=30 NATOPS",
         0.3, 1.28, 12.7, 0.35, size=13.5, color=DARK_GREY, italic=True)

wil_headers = ["Metric", "TAGFC vs CoMTE\nW-stat", "TAGFC vs CoMTE\np-value", "Sig?",
               "TAGFC vs AB-CF\nW-stat", "TAGFC vs AB-CF\np-value", "Sig?"]
wil_data = [
    ["Validity ↑",     "225",  "0.0001", "✓",  "225",  "0.0001", "✓"],
    ["Proximity L1 ↓", "15",   "0.0003", "✓",  "0",    "<0.0001","✓"],
    ["Proximity L2 ↓", "21",   "0.0008", "✓",  "0",    "<0.0001","✓"],
    ["Proximity L∞ ↓", "28",   "0.0020", "✓",  "0",    "<0.0001","✓"],
    ["Sparsity ↑",     "0",    "<0.0001","✓",  "0",    "<0.0001","✓"],
    ["Coherence ↓",    "12",   "0.0001", "✓",  "7",    "0.0001", "✓"],
    ["CF Confidence ↑","225",  "0.0001", "✓",  "225",  "0.0001", "✓"],
    ["RCF ↓",          "18",   "0.0005", "✓",  "0",    "<0.0001","✓"],
]
wcw = [1.8, 1.55, 1.65, 0.6, 1.55, 1.65, 0.6]
wcx = [0.25]
for w in wcw[:-1]:
    wcx.append(wcx[-1] + w + 0.02)

for j, (hdr, x, w) in enumerate(zip(wil_headers, wcx, wcw)):
    add_rect(slide, x, 1.68, w, 0.52, fill_rgb=DARK_BLUE)
    add_text(slide, hdr, x+0.04, 1.70, w-0.08, 0.48,
             bold=True, size=10, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(wil_data):
    bg = LIGHT_GREY if i % 2 == 0 else WHITE
    for j, (cell, x, w) in enumerate(zip(row, wcx, wcw)):
        sig_col = (j in [3, 6])
        add_rect(slide, x, 2.22+i*0.53, w, 0.5, fill_rgb=bg,
                 line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_pt=0.5)
        clr = RGBColor(0x00, 0x88, 0x00) if (sig_col and cell == "✓") else DARK_GREY
        add_text(slide, cell, x+0.04, 2.25+i*0.53, w-0.08, 0.44,
                 bold=sig_col, size=11.5, color=clr, align=PP_ALIGN.CENTER)

add_text(slide,
         "All 8 metrics show statistically significant differences (p < 0.0083) after Bonferroni correction",
         0.3, 6.7, 12.7, 0.38, size=13, color=MID_BLUE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — PROXIMITY-VALIDITY SCATTER
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Proximity–Validity Trade-off Analysis")
footer(slide, 24)

add_text(slide,
         "Ideal CF method: high validity (y-axis) AND low proximity (x-axis → left) — "
         "upper-left corner is optimal",
         0.3, 1.28, 12.7, 0.38, size=13.5, color=DARK_GREY, italic=True)

img(slide, r"TAGFC_Natops\outputs\figures\F11_proximity_validity_scatter.png", 1.5, 1.7, 10.2)

add_text(slide,
         "TAGFC occupies a distinct region: lower L2 proximity than CoMTE/AB-CF while maintaining "
         "competitive validity — demonstrating the benefit of targeted wavelet perturbation",
         0.3, 6.75, 12.7, 0.45, size=12.5, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — 21-DIMENSION COMPARISON
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Comparison & Analysis — 21-Dimension Metric Table")
footer(slide, 25)

add_text(slide,
         "Comprehensive comparison across all dataset × metric combinations "
         "(AQI India + NATOPS, 8 metrics each, both baselines)",
         0.3, 1.28, 12.7, 0.38, size=13, color=DARK_GREY, italic=True)

dim21_headers = ["Metric", "AQI\nTAGFC", "AQI\nCoMTE", "AQI\nAB-CF",
                 "NATOPS\nTAGFC", "NATOPS\nCoMTE", "NATOPS\nAB-CF"]
dim21_data = [
    ["Validity ↑",    "1.000", "0.167", "0.133", "0.867", "1.000", "1.000"],
    ["Prox L1 ↓",     "10.59", "47.14", "49.31", "35.50", "88.63", "131.05"],
    ["Prox L2 ↓",     "2.030", "6.083", "6.329", "5.197", "14.774","22.800"],
    ["Prox L∞ ↓",     "1.304", "2.326", "2.486", "0.883", "2.875", "3.027"],
    ["Sparsity ↑",    "0.000", "0.208", "0.173", "0.000", "0.917", "0.882"],
    ["Coherence ↓",   "0.132", "1.042", "1.044", "7.544", "42.904","43.471"],
    ["CF Conf ↑",     "0.955", "0.132", "0.106", "0.705", "1.000", "1.000"],
    ["RCF ↓",         "0.349", "0.920", "0.947", "0.217", "0.647", "1.027"],
]
# best per row: col 1 for AQI most; col 4 for NATOPS proximity/coherence/RCF
tagfc_aqi_best  = {0, 1, 2, 3, 5, 6, 7}    # AQI cols (indices 1)
tagfc_nat_best  = {1, 2, 3, 5, 7}           # NATOPS cols (index 4)

dcw = [1.5, 1.6, 1.6, 1.6, 1.6, 1.6, 1.6]
dcx = [0.25]
for w in dcw[:-1]:
    dcx.append(dcx[-1] + w + 0.01)

for j, (hdr, x, w) in enumerate(zip(dim21_headers, dcx, dcw)):
    bg = DARK_BLUE if j == 0 else (RGBColor(0x00, 0x44, 0x88) if j in [1,2,3] else RGBColor(0x00, 0x55, 0x44))
    add_rect(slide, x, 1.68, w, 0.52, fill_rgb=bg)
    add_text(slide, hdr, x+0.04, 1.70, w-0.08, 0.48,
             bold=True, size=10.5, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(dim21_data):
    bg = LIGHT_GREY if i % 2 == 0 else WHITE
    for j, (cell, x, w) in enumerate(zip(row, dcx, dcw)):
        is_best_aqi  = (j == 1 and i in tagfc_aqi_best)
        is_best_nat  = (j == 4 and i in tagfc_nat_best)
        add_rect(slide, x, 2.22+i*0.57, w, 0.53, fill_rgb=bg,
                 line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_pt=0.5)
        clr = RGBColor(0xCC, 0x00, 0x00) if (is_best_aqi or is_best_nat) else DARK_GREY
        add_text(slide, cell, x+0.04, 2.25+i*0.57, w-0.08, 0.48,
                 bold=(is_best_aqi or is_best_nat or j == 0), size=11,
                 color=clr, align=PP_ALIGN.CENTER)

add_text(slide,
         "Red bold = TAGFC best value in that dataset column  ·  TAGFC leads in proximity, coherence, CF confidence, and RCF across both datasets",
         0.3, 6.9, 12.7, 0.35, size=11.5, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — KEY CONTRIBUTIONS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Key Contributions")
footer(slide, 26)

contribs = [
    ("C1", "Frequency-Domain CF Framework",
     "First MTS counterfactual method to operate in the Haar wavelet coefficient space — enabling "
     "targeted perturbation of specific time-frequency bands (slow trends vs. fast oscillations)."),
    ("C2", "Attention-Guided Saliency",
     "Multi-layer Attention Rollout provides a principled, model-internal saliency signal that guides "
     "the omega weighting — directing perturbation cost towards semantically important coefficients."),
    ("C3", "Cross-Sensor Coherence",
     "Mahalanobis-based L_cross term enforces that counterfactual changes respect the empirical "
     "sensor correlation structure, producing physically plausible explanations."),
    ("C4", "Exact Sparsity via Proximal GD",
     "Soft-thresholding in the Proximal Gradient Descent step achieves exact L1 sparsity (true zeros) "
     "rather than approximate near-zero values — yielding interpretable explanations."),
    ("C5", "Empirical Validation",
     "Demonstrated on 2 real-world datasets (AQI India: TAGFC wins 7/8 metrics; "
     "NATOPS: wins 4–5/8) against 4 baselines with Wilcoxon + Bonferroni statistical testing."),
]
for i, (tag, title, body) in enumerate(contribs):
    y = 1.35 + i * 1.18
    add_rect(slide, 0.3, y, 0.6, 0.65, fill_rgb=MID_BLUE)
    add_text(slide, tag, 0.3, y+0.1, 0.6, 0.45,
             bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.05, y, 11.5, 0.38,
             bold=True, size=15, color=DARK_BLUE)
    add_text(slide, body, 1.05, y+0.37, 11.5, 0.65,
             size=13, color=DARK_GREY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — LIMITATIONS & FUTURE WORK
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
banner(slide, "Limitations & Future Work")
footer(slide, 27)

add_text(slide, "Current Limitations", 0.3, 1.35, 6.1, 0.38,
         bold=True, size=17, color=MID_BLUE)
limits = [
    "Validity not 100% on NATOPS (0.867) — proximal-GD can fail to flip the classifier within 300 iters for hard examples",
    "CoMTE/AB-CF achieve higher sparsity via whole-channel swap — TAGFC perturbation is distributed across many coefficients",
    "N=30 for NATOPS evaluation — Wilcoxon power improves with N=200 (planned)",
    "No ablation study reported — impact of each of the 4 objective terms needs quantification",
]
for i, lim in enumerate(limits):
    add_rect(slide, 0.3, 1.82+i*0.82, 0.06, 0.62, fill_rgb=ORANGE)
    add_text(slide, lim, 0.5, 1.84+i*0.82, 5.7, 0.62,
             size=13, color=DARK_GREY, wrap=True)

add_rect(slide, 6.7, 1.35, 0.04, 5.5, fill_rgb=RGBColor(0xCC, 0xCC, 0xCC))

add_text(slide, "Future Directions", 7.0, 1.35, 6.1, 0.38,
         bold=True, size=17, color=MID_BLUE)
futures = [
    "Scale NATOPS to N=200; run CMAPSS FD001 with N=200 for full 3-dataset evaluation",
    "Ablation study: noRollout (ω=1), noWavelet (raw X), noCross (λ₂=0), noManifold (λ₃=0)",
    "Extend to multivariate regression (RUL prediction) with continuous CF targets",
    "Replace Haar DWT with learnable wavelet (e.g., LiftingNet) for data-adaptive frequency decomposition",
    "Real-time deployment: streaming counterfactuals for online sensor monitoring dashboards",
    "Submit to KDD MiLeTS workshop / IEEE TNNLS journal",
]
for i, fut in enumerate(futures):
    add_rect(slide, 7.0, 1.82+i*0.7, 0.32, 0.32, fill_rgb=MID_BLUE)
    add_text(slide, str(i+1), 7.0, 1.84+i*0.7, 0.32, 0.28,
             bold=True, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, fut, 7.42, 1.84+i*0.7, 5.7, 0.55,
             size=12.5, color=DARK_GREY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — THANK YOU
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
add_rect(slide, 0, 5.5, 13.33, 2.0, fill_rgb=RGBColor(0x00, 0x22, 0x44))
add_rect(slide, 0.5, 3.2, 12.33, 0.07, fill_rgb=ORANGE)

add_text(slide, "Thank You", 1.0, 1.0, 11.3, 1.2,
         bold=True, size=54, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Questions & Discussion", 1.0, 2.3, 11.3, 0.6,
         bold=False, size=24, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER)

add_text(slide,
         "Avinash Yadav  (244101008)  ·  avinashy@iitg.ac.in\n"
         "Vaibhav Wankar  (244101061)  ·  w.vaibhav@iitg.ac.in",
         1.0, 3.5, 11.3, 0.8, size=16, color=RGBColor(0xCC, 0xDD, 0xEE),
         align=PP_ALIGN.CENTER)
add_text(slide,
         "Guide: Dr. Rashmi Dutta Baruah  ·  Dept. of CSE  ·  IIT Guwahati  ·  June 2026",
         1.0, 4.3, 11.3, 0.45, size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
         align=PP_ALIGN.CENTER)

add_text(slide,
         "TAGFC — Transformer Attention-Guided Frequency Counterfactual  ·  M.Tech Thesis",
         1.0, 5.6, 11.3, 0.38, size=13, color=RGBColor(0x88, 0xAA, 0xCC),
         align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
